#!/usr/bin/env python3
"""
extract_images.py — 从 PPT/PDF 中提取嵌入图片 / 逐页渲染 PDF 为 PNG

用法:
  # 提取嵌入图片（原有功能）
  py .scripts/extract_images.py <input_file> <output_dir> [--prefix PREFIX]

  # 逐页渲染 PDF 为 PNG（供 vision 读取）
  py .scripts/extract_images.py <input_file> <output_dir> --prefix PREFIX --pages [--dpi 150]

示例:
  # 从 PPT 提取课件图片
  py .scripts/extract_images.py "附件/L05.pptx" "附件/" --prefix CME222_L05

  # 从 PDF 提取嵌入图片
  py .scripts/extract_images.py "附件/L01.pdf" "附件/" --prefix CME213_L01

  # PDF 逐页渲染为 PNG（用于 Claude Code vision 读取）
  py .scripts/extract_images.py "附件/L01.pdf" "附件/_pages/" --prefix CME213_L01 --pages --dpi 150

输出:
  - 提取模式：图片文件（去重+压缩后）+ 控制台报告
  - 渲染模式：每页一个 PNG 文件 + 控制台报告

依赖:
  pip install PyMuPDF Pillow python-pptx
"""

import argparse
import hashlib
import io
import sys
from pathlib import Path

# 修复 Windows 终端编码
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")

# 图片大小阈值
MIN_SIZE_BYTES = 5 * 1024       # <5KB 视为装饰图，跳过
MAX_SIZE_BYTES = 300 * 1024     # >300KB 需要压缩
MAX_DIMENSION = 1200            # 压缩后最大宽度/高度


def md5_hash(data: bytes) -> str:
    """计算 MD5 哈希值"""
    return hashlib.md5(data).hexdigest()


def is_mostly_solid(image_data: bytes) -> bool:
    """检测图片是否为纯色/接近纯色（90% 以上像素相同）"""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(image_data))
        # 缩小到 50x50 再检测，提高速度
        img_small = img.resize((50, 50))
        pixels = list(img_small.getdata())
        if not pixels:
            return True
        # 统计最常见颜色的占比
        from collections import Counter
        most_common_count = Counter(pixels).most_common(1)[0][1]
        return most_common_count / len(pixels) > 0.9
    except Exception:
        return False


def compress_image(image_data: bytes, ext: str) -> bytes:
    """压缩图片到合理大小"""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(image_data))

        # 按比例缩放
        w, h = img.size
        if max(w, h) > MAX_DIMENSION:
            ratio = MAX_DIMENSION / max(w, h)
            new_size = (int(w * ratio), int(h * ratio))
            img = img.resize(new_size, Image.LANCZOS)

        # 保存到 bytes
        buf = io.BytesIO()
        if ext.lower() in (".jpg", ".jpeg"):
            # 转换 RGBA → RGB（JPEG 不支持透明通道）
            if img.mode == "RGBA":
                img = img.convert("RGB")
            img.save(buf, format="JPEG", quality=85, optimize=True)
        else:
            img.save(buf, format="PNG", optimize=True)
        return buf.getvalue()
    except Exception as e:
        print(f"  压缩失败: {e}", file=sys.stderr)
        return image_data


def extract_from_pptx(pptx_path: Path) -> list[tuple[str, bytes]]:
    """从 PPTX 提取图片，返回 [(slide_info, image_bytes), ...]"""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print("需要安装 python-pptx: pip install python-pptx", file=sys.stderr)
        return []

    images = []
    prs = Presentation(str(pptx_path))

    for slide_idx, slide in enumerate(prs.slides, 1):
        fig_count = 0
        for shape in slide.shapes:
            # shape_type 13 = Picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                fig_count += 1
                blob = shape.image.blob
                ext = shape.image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                slide_info = f"s{slide_idx:02d}_fig{fig_count:02d}.{ext}"
                images.append((slide_info, blob))

    return images


def extract_from_pdf(pdf_path: Path) -> list[tuple[str, bytes]]:
    """从 PDF 提取图片，返回 [(page_info, image_bytes), ...]"""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("需要安装 PyMuPDF: pip install PyMuPDF", file=sys.stderr)
        return []

    images = []
    doc = fitz.open(str(pdf_path))

    for page_idx, page in enumerate(doc, 1):
        img_list = page.get_images(full=True)
        for img_idx, img_info in enumerate(img_list, 1):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                image_data = base_image["image"]
                ext = base_image["ext"]
                page_info = f"p{page_idx:02d}_fig{img_idx:02d}.{ext}"
                images.append((page_info, image_data))
            except Exception as e:
                print(f"  页 {page_idx} 图 {img_idx} 提取失败: {e}", file=sys.stderr)

    doc.close()
    return images


def process_images(
    raw_images: list[tuple[str, bytes]],
    output_dir: Path,
    prefix: str,
) -> dict:
    """处理提取的图片：去重、过滤、压缩、保存"""
    output_dir.mkdir(parents=True, exist_ok=True)

    seen_hashes: set[str] = set()
    stats = {"extracted": 0, "skipped_small": 0, "skipped_solid": 0,
             "skipped_dup": 0, "compressed": 0, "saved": 0}

    for info, data in raw_images:
        stats["extracted"] += 1

        # 跳过太小的图
        if len(data) < MIN_SIZE_BYTES:
            stats["skipped_small"] += 1
            continue

        # 跳过纯色图
        if is_mostly_solid(data):
            stats["skipped_solid"] += 1
            continue

        # 去重
        h = md5_hash(data)
        if h in seen_hashes:
            stats["skipped_dup"] += 1
            continue
        seen_hashes.add(h)

        # 压缩
        ext = Path(info).suffix
        if len(data) > MAX_SIZE_BYTES:
            data = compress_image(data, ext)
            stats["compressed"] += 1

        # 保存
        filename = f"{prefix}_{info}"
        out_path = output_dir / filename
        out_path.write_bytes(data)
        stats["saved"] += 1
        print(f"  保存: {filename} ({len(data) / 1024:.0f} KB)")

    return stats


def render_pages_to_png(pdf_path: Path, output_dir: Path, prefix: str, dpi: int = 150) -> int:
    """把 PDF 每页渲染成 PNG，供 Claude Code vision 读取。返回渲染页数。"""
    try:
        import fitz
    except ImportError:
        print("需要安装 PyMuPDF: pip install PyMuPDF", file=sys.stderr)
        return 0

    doc = fitz.open(str(pdf_path))
    output_dir.mkdir(parents=True, exist_ok=True)
    total = len(doc)

    print(f"逐页渲染: {pdf_path.name} ({total} 页, {dpi} DPI)")

    for page_idx, page in enumerate(doc, 1):
        pix = page.get_pixmap(dpi=dpi)
        filename = f"{prefix}_page{page_idx:02d}.png"
        out_path = output_dir / filename
        pix.save(str(out_path))
        size_kb = out_path.stat().st_size / 1024
        print(f"  渲染: {filename} ({size_kb:.0f} KB)")

    doc.close()
    print(f"\n共渲染 {total} 页到 {output_dir}")
    return total


def main():
    parser = argparse.ArgumentParser(
        description="从 PPT/PDF 提取嵌入图片，或逐页渲染 PDF 为 PNG"
    )
    parser.add_argument("input_file", help="PPT 或 PDF 文件路径")
    parser.add_argument("output_dir", help="图片输出目录")
    parser.add_argument(
        "--prefix", default="img",
        help="输出文件名前缀（如 CME222_L05）"
    )
    # 逐页渲染模式
    parser.add_argument(
        "--pages", action="store_true",
        help="逐页渲染 PDF 为 PNG（供 vision 读取），而非提取嵌入图片"
    )
    parser.add_argument(
        "--dpi", type=int, default=150,
        help="逐页渲染的 DPI（默认 150，越高越清晰但文件越大）"
    )
    args = parser.parse_args()

    input_path = Path(args.input_file)
    output_dir = Path(args.output_dir)

    if not input_path.exists():
        print(f"文件不存在: {input_path}", file=sys.stderr)
        return 1

    ext = input_path.suffix.lower()

    # --pages 模式：逐页渲染 PDF
    if args.pages:
        if ext != ".pdf":
            print(f"--pages 模式仅支持 PDF 文件，当前: {ext}", file=sys.stderr)
            return 1
        count = render_pages_to_png(input_path, output_dir, args.prefix, args.dpi)
        return 0 if count > 0 else 1

    # 默认模式：提取嵌入图片
    if ext in (".pptx", ".ppt"):
        print(f"从 PPT 提取图片: {input_path.name}")
        raw_images = extract_from_pptx(input_path)
    elif ext == ".pdf":
        print(f"从 PDF 提取图片: {input_path.name}")
        raw_images = extract_from_pdf(input_path)
    else:
        print(f"不支持的格式: {ext}（支持 .pptx, .pdf）", file=sys.stderr)
        return 1

    if not raw_images:
        print("未找到嵌入图片")
        return 0

    print(f"找到 {len(raw_images)} 张原始图片，开始处理...")
    stats = process_images(raw_images, output_dir, args.prefix)

    # 输出报告
    print(f"\n--- 提取报告 ---")
    print(f"原始图片: {stats['extracted']}")
    print(f"跳过（太小）: {stats['skipped_small']}")
    print(f"跳过（纯色）: {stats['skipped_solid']}")
    print(f"跳过（重复）: {stats['skipped_dup']}")
    print(f"压缩: {stats['compressed']}")
    print(f"最终保存: {stats['saved']}")

    return 0


if __name__ == "__main__":
    sys.exit(main())
